import pathlib
from pypdf import PdfReader
from pptx import Presentation

root = pathlib.Path(__file__).resolve().parent
out_dir = root / "extracted"
out_dir.mkdir(exist_ok=True)

pdfs = [
    "IT22069122.pdf",
    "IT22088000.pdf",
    "IT22360496.pdf",
    "IT22564986.pdf",
    "TAF_25-26J-010.pdf",
]
pptx_files = ["Proposal Presentation_25-26J-010.pptx"]

def save_text(src: pathlib.Path, text: str) -> None:
    text = text.replace("\r\n", "\n")
    (out_dir / f"{src.stem}.txt").write_text(text, encoding="utf-8")

def extract_pdfs():
    for name in pdfs:
        p = root / name
        try:
            reader = PdfReader(str(p))
            parts = []
            for page in reader.pages:
                parts.append(page.extract_text() or "")
            save_text(p, "\n\n".join(parts))
            print(f"OK PDF {name}")
        except Exception as e:
            print(f"FAIL PDF {name}: {e}")

def extract_pptx():
    for name in pptx_files:
        p = root / name
        try:
            pres = Presentation(str(p))
            lines = []
            for i, slide in enumerate(pres.slides, 1):
                lines.append(f"--- Slide {i} ---")
                for shape in slide.shapes:
                    if not hasattr(shape, "text"):
                        continue
                    txt = shape.text.strip()
                    if txt:
                        lines.append(txt)
            save_text(p, "\n".join(lines))
            print(f"OK PPTX {name}")
        except Exception as e:
            print(f"FAIL PPTX {name}: {e}")

def main():
    extract_pdfs()
    extract_pptx()

if __name__ == "__main__":
    main()
